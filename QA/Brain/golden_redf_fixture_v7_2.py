from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from delivery_test_fixtures_v7_2 import bad_pptx

def add_text(sl,text,x,y,w,h,size=24):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); p=tb.text_frame.paragraphs[0]; p.text=text; p.font.size=Pt(size); return tb

def add_table(sl,rows,cols,x,y,w,h,data):
    t=sl.shapes.add_table(rows,cols,Inches(x),Inches(y),Inches(w),Inches(h)).table
    for r,row in enumerate(data):
        for c,v in enumerate(row): t.cell(r,c).text=str(v)
    return t

def build(path,hero):
    prs=Presentation(); prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    s=prs.slides.add_slide(blank); s.shapes.add_picture(str(hero),Inches(.6),Inches(.8),width=Inches(6)); add_text(s,'مشروع الابتكار والثقافة الرقمية',7,1.5,5.4,1.5,30)
    s=prs.slides.add_slide(blank); add_text(s,'HOLD — فرصة جذابة، لكن الجاهزية غير مثبتة',1,1,11,1.2,30); add_table(s,3,3,1,2.6,11,3,[['Why attractive','Why not GO','What changes decision'],['Transformation','Evidence gaps','Close 8 gates'],['Strategic relevance','Commercial unknowns','Bid readiness sprint']])
    s=prs.slides.add_slide(blank); add_text(s,'30',.8,1.3,2.5,1.4,54);add_text(s,'POCs',.8,2.6,2.5,.8,20);add_text(s,'20',4,1.3,2.5,1.4,54);add_text(s,'Services',4,2.6,2.5,.8,20);add_text(s,'24',7.2,1.3,2.5,1.4,54);add_text(s,'Resources',7.2,2.6,2.5,.8,20);add_text(s,'12',10.3,1.3,2,1.4,54);add_text(s,'Months support',10.1,2.6,2.3,.8,18)
    s=prs.slides.add_slide(blank); add_text(s,'Strategy → Challenges → Ideas → POCs → Services',1,1,11,1,28); add_text(s,'Platform + AI + Measurement form the operating backbone',2,3,9,1.4,24)
    s=prs.slides.add_slide(blank); add_table(s,5,3,.8,1.2,11.8,5,[['Issue','Why it matters','Action'],['Open-source contradiction','Architecture / licensing','Clarify'],['Duration incomplete','Pricing','Clarify'],['Missing annex','Compliance','Obtain'],['Template noise','Contract scope','Confirm']])
    s=prs.slides.add_slide(blank)
    for txt,y in [('EXPERIENCE / SERVICES',1.2),('PLATFORM / INTEGRATION',2.3),('DATA / AI',3.4),('SECURITY / INFRASTRUCTURE',4.5)]: add_text(s,txt,2,y,9,.7,22)
    s=prs.slides.add_slide(blank); add_table(s,6,4,.7,1,12,5.5,[['BOQ','Qty','Hidden effort','Pricing sensitivity'],['Strategy',1,'Medium','Medium'],['Program',1,'High','High'],['Services',20,'Very High','Very High'],['POCs',30,'Extreme','Extreme'],['Support',12,'High','High']])
    s=prs.slides.add_slide(blank); add_text(s,'Experience',1,1,3,.7,20);add_text(s,'Integration / APIs',1,2,3,.7,20);add_text(s,'Data & AI',1,3,3,.7,20);add_text(s,'Security',1,4,3,.7,20);add_text(s,'On-Prem Infrastructure',1,5,3,.7,20);add_text(s,'Critical contradiction:\nRequired stack ↔ open-source restriction',6,2.3,5.5,2,26)
    s=prs.slides.add_slide(blank); add_text(s,'24',1,1.5,3,2,64);add_text(s,'mandatory named resources',1,3.3,4,1,22);add_table(s,4,2,6,1.5,5.5,4,[['Dimension','Commitment'],['Roles','12'],['Interviews','Required'],['Replacement','≤2 weeks']])
    s=prs.slides.add_slide(blank); cd=ChartData();cd.categories=['Experience','Solution','Team','Proposal'];cd.add_series('Points',(20,30,20,30));s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(.8),Inches(1.2),Inches(7),Inches(4.8),cd);add_text(s,'87.5%',8.5,1.6,3,1.5,46);add_text(s,'of remaining points needed if experience = 0/20',8.2,3.2,4,1.5,20)
    s=prs.slides.add_slide(blank); vals=[('1%','Bid bond'),('5%','Performance guarantee'),('20%','Penalty ceiling'),('24','Resources'),('12','Support months')]
    for i,(v,l) in enumerate(vals): x=.6+i*2.5;add_text(s,v,x,1.8,2,1.2,38);add_text(s,l,x,3,2,1,16)
    s=prs.slides.add_slide(blank); add_text(s,'Impact × Urgency',.8,.6,4,.8,26)
    for i,t in enumerate(['Contract duration','Open-source contradiction','30 POC definition','Payment schedule']): add_text(s,t,1+(i%2)*6,1.8+(i//2)*2.2,5,1.2,20)
    s=prs.slides.add_slide(blank); add_table(s,9,2,2,1,9,5.7,[['Gate','Status'],['Saudi experience','Unknown'],['ISO','Unknown'],['Team capacity','Unknown'],['Tech contradiction','Open'],['BOQ decomposition','Open'],['Commercial model','Open'],['Score path','Pending'],['Price floor','Pending']])
    s=prs.slides.add_slide(blank); add_text(s,'Bid Readiness Sprint',1,1,11,1,30)
    for i,t in enumerate(['Evidence','Team','Clarifications','Cost model','Technical score','Price floor']): add_text(s,t,.6+i*2.05,3,1.7,1,17)
    prs.save(path)
