#!/usr/bin/env python3
from __future__ import annotations
from pathlib import Path
import sys,json,tempfile
ROOT=Path(__file__).resolve().parents[2]
RUNTIME=ROOT/'Rashad/Brain/runtime'
sys.path.insert(0,str(RUNTIME)); sys.path.insert(0,str(ROOT/'QA/Brain'))
from brain.expert_router import route_experts
from brain.expert_runtime import execute_expert_council
from brain.provider import ScriptedTestProvider,NoExecutionProvider
from brain.orchestrator import run_brain
from brain.artifact_council_runtime import execute_artifact_councils
from brain.artifact_brain import generate_communication_hypotheses, route_artifact_councils
from artifact_delivery_orchestrator import promote_page_to_production
from brain.product_inspector import inspect_pptx
from orchestrator import route as qa_route
from pptx import Presentation
from pptx.util import Inches,Pt

checks=[]
def ck(name,cond,detail=None): checks.append({'name':name,'status':'PASS' if cond else 'FAIL','detail':detail})

# 1. Big Brain routing must select meaningful executable expertise, not just register names.
technical={'task_id':'T-AI','rfp_role':'TECHNICAL_REQUIREMENTS','critical':True,'rendered':True,'question':'Design AI platform with APIs, data integration, cybersecurity, cloud, Docker Kubernetes, backend frontend mobile UX accessibility'}
tr=route_experts(technical); ts=set(tr['selected_experts'])
ck('technical_route_pass',tr['status']=='PASS',tr)
for rid in ['SME-ENTERPRISE-ARCH','SME-SOLUTION-ARCH','SIM-CIO-CTO','GOV-TECHNICAL-FEASIBILITY','SME-AI-ENGINEERING','SME-CYBER','SME-CLOUD-INFRA']:
    ck('technical_has_'+rid,rid in ts,sorted(ts))
ck('technical_bounded',tr['selected_count']<=tr['max_active_experts'],tr['selected_count'])

commercial={'task_id':'T-COM','rfp_role':'COMMERCIAL_EXPOSURE','critical':True,'question':'fixed price margin payment guarantees penalties working capital contract risk'}
cr=route_experts(commercial); cs=set(cr['selected_experts'])
for rid in ['SME-PRICING','SME-COMMERCIAL','SME-LEGAL-CONTRACT','SIM-CFO','GOV-FINANCIAL-TRUTH']:
    ck('commercial_has_'+rid,rid in cs,sorted(cs))

# 2. JSON key names must never trigger domain routing.
key_attack={'task_id':'KEY','rfp_role':'MANAGEMENT_DECISION','critical':True,'rfp_procurement_cyber_cloud_sap':'plain strategy only','payload':{'fake_ai_key':'ordinary management question'}}
ka=route_experts(key_attack)
ck('semantic_values_not_keys_no_sap','SAP' not in ka.get('matched_domains',[]),ka)
ck('semantic_values_not_keys_no_cyber','CYBER_PRIVACY' not in ka.get('matched_domains',[]),ka)
# value trigger still works
kv=route_experts({'task_id':'VAL','rfp_role':'TECHNICAL_REQUIREMENTS','critical':True,'payload':'SAP integration and cybersecurity'})
ck('semantic_values_still_route_sap','SAP' in kv.get('matched_domains',[]),kv)
ck('semantic_values_still_route_cyber','CYBER_PRIVACY' in kv.get('matched_domains',[]),kv)

# 3. Selected experts must actually execute in isolated invocations.
ledger=execute_expert_council(technical,{'answer_first_thesis':'Integrated AI architecture'},ScriptedTestProvider())
ck('expert_ledger_pass',ledger['status']=='PASS',ledger.get('errors'))
ck('expert_required_equals_executed',set(ledger['required_experts'])<=set(ledger['executed_experts']),{'required':ledger['required_experts'],'executed':ledger['executed_experts']})
ck('expert_isolated_actor_ids',len({x.get('actor_id') for x in ledger['invocations']})==len(ledger['invocations']))
ck('expert_isolated_contexts',len({x.get('isolated_context_id') for x in ledger['invocations']})==len(ledger['invocations']))
blocked_ledger=execute_expert_council(technical,{'answer_first_thesis':'x'},NoExecutionProvider())
ck('registered_but_not_executed_blocks',blocked_ledger['status']=='BLOCKED' and bool(blocked_ledger['errors']),blocked_ledger.get('errors'))

# 4. Full Brain cannot reach cognitive lock without executable expertise.
bs=run_brain({'task_id':'T-AI','rfp_role':'TECHNICAL_REQUIREMENTS','critical':True,'rendered':True,'question':technical['question'],'evidence':[{'id':'E1'}]},ScriptedTestProvider())
ck('brain_reaches_cognitive_lock_with_execution',bs.get('state')=='COGNITIVE_LOCKED',bs.get('state'))
ck('brain_contains_expert_execution_ledger',(bs.get('expert_execution_ledger') or {}).get('status')=='PASS')
ck('brain_coverage_pass',(bs.get('coverage') or {}).get('status')=='PASS',bs.get('coverage'))

# 5. Artifact Brain councils are executable and Arabic routing is explicit.
graph={'nodes':[{'id':'N1','label':'Evidence'},{'id':'N2','label':'Decision'}],'edges':[{'source':'N1','target':'N2','relation':'SUPPORTS'}]}
content={'thesis':'القرار يحتاج إثبات الجاهزية قبل التقديم','answer_first_thesis':'القرار يحتاج إثبات الجاهزية قبل التقديم','evidence':['E1'],'numbers':[24,30,20],'decision':'HOLD','rfp':'منافسة حكومية'}
ar=route_artifact_councils(graph,content,'AR','PRE_CONCEPT')
ck('artifact_route_pass',ar['status']=='PASS',ar)
ck('artifact_arabic_rtl_council_routed','ARABIC_RTL_COUNCIL' in ar['active_councils'],ar['active_councils'])
ck('artifact_meta_council_routed','META_EXPERTISE_ROUTER' in ar['active_councils'],ar['active_councils'])
ac=execute_artifact_councils(graph,content,ScriptedTestProvider(),'AR','PRE_CONCEPT')
ck('artifact_council_execution_pass',ac['status']=='PASS',ac.get('errors'))
ck('artifact_all_required_councils_executed',set(x['council_id'] for x in ac['required_executions'])<=set(ac['executed_councils']),ac.get('executed_councils'))
ac_no=execute_artifact_councils(graph,content,NoExecutionProvider(),'AR','PRE_CONCEPT')
ck('artifact_metadata_without_execution_blocks',ac_no['status']=='BLOCKED',ac_no.get('errors'))
h=generate_communication_hypotheses(graph,content,5)
ck('five_unique_communication_strategies',h['status']=='PASS' and h['distinct_communication_strategies']==5,h)
ck('artifact_not_diagram_only',h['diagram_only_search'] is False,h)
ck('artifact_has_minimal_candidate',h['contains_minimal_hypothesis'] is True,h)

# 6. Direct production without upstream Brain/council proof must fail before renderer.
def never_renderer(*a,**k): raise AssertionError('renderer must not be called')
def never_reviewer(*a,**k): raise AssertionError('reviewer must not be called')
prom=promote_page_to_production({'status':'PASS'},never_renderer,never_reviewer,{'material_claim_ids':['C1']},graph,tempfile.mkdtemp())
ck('direct_production_without_brain_blocks',prom.get('status')=='BLOCKED' and prom.get('reason')=='BRAIN_EXECUTION_PROOF_NOT_PROVEN',prom)


# Fake status-only proof must not pass structural proof validation.
fake={'state':'COGNITIVE_LOCKED','coverage':{'status':'PASS'},'expert_execution_ledger':{'status':'PASS'}}
fake_prom=promote_page_to_production({'artifact_council_execution':{'status':'PASS'},'art_direction_execution':{'status':'PASS'}},never_renderer,never_reviewer,{'material_claim_ids':[]},graph,tempfile.mkdtemp(),brain_session=fake,page_contract={'page_id':'P'},artifact_intent={'type':'X'},evidence_lineage=[{'evidence':'E1'}],brand_preflight={'status':'PASS','rubix_asset_status':'VERIFIED'})
ck('fake_status_only_brain_proof_blocked',fake_prom.get('status')=='BLOCKED' and fake_prom.get('reason')=='BRAIN_EXECUTION_PROOF_NOT_PROVEN',fake_prom)

# 7. Product inspector must block the exact failure pattern: repetitive equal cards.
def card_deck(path):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    # cover
    s=prs.slides.add_slide(blank); tb=s.shapes.add_textbox(Inches(1),Inches(2),Inches(11),Inches(1)); tb.text_frame.text='ملخص المنافسة'
    for ix in range(8):
        s=prs.slides.add_slide(blank)
        for r in range(2):
            for c in range(3):
                tb=s.shapes.add_textbox(Inches(.7+c*4.1),Inches(1.1+r*2.7),Inches(3.4),Inches(1.9)); tb.text_frame.text=f'بطاقة {ix+1}-{r*3+c+1}'
                tb.text_frame.paragraphs[0].font.size=Pt(18)
    prs.save(path)
with tempfile.TemporaryDirectory() as td:
    p=Path(td)/'cards.pptx'; card_deck(p); pi=inspect_pptx(p)
    ck('equal_card_grid_deck_blocked',pi['status']=='BLOCKED',pi.get('blockers'))
    ck('equal_card_grid_specific_blocker',bool({'PPTX_EQUAL_CARD_GRID_OVERUSE','EQUAL_CARD_GRID_OVERUSE'} & set(pi.get('blockers',[]))),pi.get('blockers'))

# 8. User-visible Arabic QA must route actual-pixel/simplicity/delivery plus Arabic council.
qr=qa_route({'rendered':True,'user_visible':True,'deck_level':True,'language':'AR'})
for qid in ['Q05_ARABIC_RTL_TYPOGRAPHY','Q11_ACTUAL_PIXEL_PRODUCT_REVIEW','Q12_EXECUTIVE_SIMPLICITY_ARTIFACT_SKEPTIC','Q13_DELIVERY_INTEGRITY_REPAIR_CLOSURE']:
    ck('qa_routes_'+qid,qid in qr,qr)

out={'suite':'Rashad v7.2 Brain Coherence & Executable Intelligence Audit','status':'PASS' if all(x['status']=='PASS' for x in checks) else 'FAIL','passed':sum(x['status']=='PASS' for x in checks),'total':len(checks),'checks':checks}
path=ROOT/'QA/Certification/V7_2_BRAIN_COHERENCE_AUDIT.json'; path.parent.mkdir(exist_ok=True); path.write_text(json.dumps(out,ensure_ascii=False,indent=2),encoding='utf-8')
print(json.dumps({'suite':out['suite'],'status':out['status'],'passed':out['passed'],'total':out['total'],'failures':[x for x in checks if x['status']!='PASS']},ensure_ascii=False,indent=2))
raise SystemExit(0 if out['status']=='PASS' else 1)
