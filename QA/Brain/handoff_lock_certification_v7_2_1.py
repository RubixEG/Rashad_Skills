#!/usr/bin/env python3
from __future__ import annotations
from pathlib import Path
import copy, hashlib, json, sys, tempfile
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw
ROOT=Path(__file__).resolve().parents[2]
sys.path.insert(0,str(ROOT/'Rashad/Brain/runtime')); sys.path.insert(0,str(Path(__file__).parent))
from brain.exact_handoff import verify_exact_artifact_handoff, issue_exact_handoff_certificate
from brain.product_inspector import inspect_pptx, sha256_file
from production_test_fixtures_v7_3 import build_production_projection
OUT=ROOT/'QA/Certification/HANDOFF_LOCK_CERTIFICATION_V7_2_1.json'

def image(path):
    im=Image.new('RGB',(640,360),'white'); d=ImageDraw.Draw(im); d.rectangle((80,70,560,290),outline='black',width=6); d.text((200,160),'VISUAL CONCEPT',fill='black'); im.save(path)

def deck(path,img):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    s=prs.slides.add_slide(blank); s.shapes.add_picture(str(img),Inches(1),Inches(1.4),width=Inches(5.5)); s.shapes.add_textbox(Inches(7),Inches(1.5),Inches(5),Inches(1)).text_frame.text='Image-led page'
    s=prs.slides.add_slide(blank); s.shapes.add_table(2,2,Inches(1),Inches(1.7),Inches(7),Inches(2)); s.shapes.add_textbox(Inches(1),Inches(.6),Inches(10),Inches(.7)).text_frame.text='Evidence table'
    s=prs.slides.add_slide(blank); s.shapes.add_textbox(Inches(1),Inches(1.8),Inches(10),Inches(1.4)).text_frame.text='24'; s.shapes.add_textbox(Inches(1),Inches(3.5),Inches(10),Inches(.8)).text_frame.text='Hero metric'
    prs.save(path)

def dossier_for(ppt,projection):
    h=sha256_file(ppt)
    pages=[]
    for i,st in enumerate(['IMAGE_LED','TABLE_LED','NUMBER_LED'],1):
        proof=dict(projection['pages'][i-1]); render=projection['renders'][i-1]
        proof.update({
          'page_id':f'P{i:02d}','selected_strategy':st,'render_kind':'PRODUCTION_PAGE_RENDER',
          'production_render_id':render['production_render_id'],'selected_render_hash':render['actual_render_hash'],
          'actual_pixel_review':{'status':'PASS','actual_render_hash':render['actual_render_hash']},
          'hero_metric_proven':st=='NUMBER_LED','structured_grid_rendered':st=='TABLE_LED'
        })
        pages.append(proof)
    product=inspect_pptx(ppt,pages)
    return {'schema':'RASHAD_USER_VISIBLE_ARTIFACT_DELIVERY_V2','classification':'USER_VISIBLE_ARTIFACT_DRAFT','output_file_sha256':h,'pages':pages,'deck_pixel_review':{'status':'PASS','independent':True,'deck_sha256':h},'product_inspection':product}

def add(rows,name,ok,detail=None): rows.append({'name':name,'status':'PASS' if ok else 'FAIL','detail':detail})

def main():
    rows=[]
    with tempfile.TemporaryDirectory() as td:
        td=Path(td); projection=build_production_projection(td/'production',['IMAGE_LED','TABLE_LED','NUMBER_LED']); assert projection.get('status')=='PASS', projection; ppt=Path(projection['pptx_path']); d=dossier_for(ppt,projection)
        v=verify_exact_artifact_handoff(ppt,d); c=issue_exact_handoff_certificate(ppt,d)
        add(rows,'matching_exact_file_certified',v['status']=='HANDOFF_ALLOWED' and c['status']=='CERTIFIED_FOR_HANDOFF',v.get('blockers'))
        x=copy.deepcopy(d); x['output_file_sha256']='0'*64; r=verify_exact_artifact_handoff(ppt,x); add(rows,'dossier_sha_mismatch_blocks',r['status']=='BLOCK_HANDOFF' and 'DELIVERED_PPTX_SHA_MISMATCH_DOSSIER' in r['blockers'],r['blockers'])
        x=copy.deepcopy(d); x['deck_pixel_review']['deck_sha256']='1'*64; r=verify_exact_artifact_handoff(ppt,x); add(rows,'deck_review_sha_mismatch_blocks','DELIVERED_PPTX_SHA_MISMATCH_DECK_PIXEL_REVIEW' in r['blockers'],r['blockers'])
        x=copy.deepcopy(d); x['product_inspection']['pptx_sha256']='2'*64; r=verify_exact_artifact_handoff(ppt,x); add(rows,'product_sha_mismatch_blocks','DELIVERED_PPTX_SHA_MISMATCH_PRODUCT_INSPECTION' in r['blockers'],r['blockers'])
        x=copy.deepcopy(d); x['pages']=x['pages'][:-1]; r=verify_exact_artifact_handoff(ppt,x); add(rows,'slide_count_mismatch_blocks','DELIVERED_SLIDE_COUNT_MISMATCH_DOSSIER_PAGES' in r['blockers'],r['blockers'])
        x=copy.deepcopy(d); x['pages'][1]['actual_pixel_review']={'status':'NOT_EXECUTED'}; r=verify_exact_artifact_handoff(ppt,x); add(rows,'pixel_count_mismatch_blocks','PIXEL_REVIEW_COUNT_MISMATCH_DELIVERED_SLIDES' in r['blockers'],r['blockers'])
        x=copy.deepcopy(d); x['pages'][1]['production_render_id']=None; r=verify_exact_artifact_handoff(ppt,x); add(rows,'production_render_count_mismatch_blocks','PRODUCTION_RENDER_COUNT_MISMATCH_DELIVERED_SLIDES' in r['blockers'],r['blockers'])
        cert=issue_exact_handoff_certificate(ppt,d); cert['pptx_sha256']='3'*64; r=verify_exact_artifact_handoff(ppt,d,certificate=cert); add(rows,'certificate_reuse_wrong_bytes_blocks','HANDOFF_CERTIFICATE_PPTX_SHA_MISMATCH' in r['blockers'],r['blockers'])
        trace=td/'trace.md'; trace.write_text('24-page PPTX built\nActual pixel reviews: 24/24\n',encoding='utf-8'); r=verify_exact_artifact_handoff(ppt,d,trace_path=trace); add(rows,'trace_different_page_count_blocks','FINAL_TRACE_DESCRIBES_DIFFERENT_PAGE_COUNT_THAN_DELIVERED_FILE' in r['blockers'],r['blockers'])
    status='PASS' if all(x['status']=='PASS' for x in rows) else 'FAIL'; out={'suite':'Rashad v7.2.1 Exact Artifact Handoff Lock Certification','status':status,'passed':sum(x['status']=='PASS' for x in rows),'total':len(rows),'results':rows}; OUT.write_text(json.dumps(out,ensure_ascii=False,indent=2),encoding='utf-8'); print(json.dumps({'suite':out['suite'],'status':status,'passed':out['passed'],'total':out['total'],'failed':[x for x in rows if x['status']!='PASS']},ensure_ascii=False,indent=2)); return 0 if status=='PASS' else 2
if __name__=='__main__': raise SystemExit(main())
