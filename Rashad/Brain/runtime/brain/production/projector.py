from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import hashlib

def build_image_master_pptx(page_images,out_path):
    prs=Presentation(); prs.slide_width=Inches(13.333333); prs.slide_height=Inches(7.5)
    # remove initial slide layouts only used as templates; start with blank slides
    blank=prs.slide_layouts[6]
    for img in page_images:
        s=prs.slides.add_slide(blank); s.shapes.add_picture(str(img),0,0,width=prs.slide_width,height=prs.slide_height)
    # python-pptx creates no initial slide; save
    out=Path(out_path); out.parent.mkdir(parents=True,exist_ok=True); prs.save(out)
    return {'status':'PASS','pptx_path':str(out),'pptx_sha256':hashlib.sha256(out.read_bytes()).hexdigest(),'slide_count':len(page_images),'projection_mode':'RASTER_PROJECTION_FROM_SEMANTIC_HTML_MASTER'}
