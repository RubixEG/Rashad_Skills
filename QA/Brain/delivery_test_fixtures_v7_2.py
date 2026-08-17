from __future__ import annotations
from pathlib import Path
import hashlib
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image,ImageDraw

# Neutral v7.2 test fixtures. This module contains no authority logic; it only builds
# deterministic files/records consumed by current v7.2 certification and Red Team suites.
def sha256_file(path):
    h=hashlib.sha256()
    with open(path,'rb') as f:
        for chunk in iter(lambda:f.read(1024*1024),b''): h.update(chunk)
    return h.hexdigest()

def png(path,text):
    im=Image.new('RGB',(800,450),'white'); d=ImageDraw.Draw(im); d.text((40,40),text,fill='black'); d.rectangle((80,120,720,360),outline='black',width=3); im.save(path); return sha256_file(path)

def review(h,score=92,ind=True):
    dims=['message_clarity','five_second_comprehension','visual_form_fitness','simplicity','executive_hierarchy','evidence_legibility','artifact_usefulness','specificity_to_page','rtl_typography','brand_fidelity','production_quality']
    return {'status':'PASS','independent':ind,'review_id':'REV-'+h[:10],'actor_id':'QA-INDEPENDENT-PIXEL','actual_render_hash':h,'scores':{k:score for k in dims},'generic_layout_swap_test':'PASS','artifact_skeptic_test':'PASS','five_second_test':'PASS','hard_blockers':[],'artifact_truth_score':score,'ceqs_score':score}

def hyps():
    ss=[('STATEMENT_LED','MINIMAL'),('NUMBER_LED','MINIMAL'),('EVIDENCE_LED','ANALYTICAL'),('COMPARISON_LED','ANALYTICAL'),('HYBRID_EXHIBIT','HYBRID')]
    return [{'id':f'H{i+1}','communication_strategy':s,'strategy_family':f,'structural_signature':f'{s}|SIG{i}'} for i,(s,f) in enumerate(ss)]

def good_pptx(path,imgpath):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    sl=prs.slides.add_slide(blank); sl.shapes.add_picture(str(imgpath),Inches(.6),Inches(.7),width=Inches(6.1)); tb=sl.shapes.add_textbox(Inches(7.2),Inches(1.4),Inches(5.3),Inches(2)); tb.text_frame.text='مشروع الابتكار والثقافة الرقمية'
    sl=prs.slides.add_slide(blank); tb=sl.shapes.add_textbox(Inches(1),Inches(2.2),Inches(11),Inches(2)); p=tb.text_frame.paragraphs[0]; p.text='الفرصة جذابة، لكن قرار الدخول يحتاج إغلاق بوابات الجاهزية'; p.font.size=Pt(28)
    sl=prs.slides.add_slide(blank); tb=sl.shapes.add_textbox(Inches(1),Inches(1.4),Inches(5),Inches(2.2)); p=tb.text_frame.paragraphs[0]; p.text='87.5%'; p.font.size=Pt(64); tb2=sl.shapes.add_textbox(Inches(6.3),Inches(2),Inches(5),Inches(2)); tb2.text_frame.text='المطلوب من بقية النقاط إذا فقدنا نقاط الخبرة'
    sl=prs.slides.add_slide(blank); tbl=sl.shapes.add_table(4,4,Inches(.8),Inches(1.5),Inches(11.7),Inches(4.5)).table
    vals=[['Issue','Evidence','Impact','Action'],['مدة العقد','غير مكتملة','تسعير','Clarify'],['Open Source','تعارض','معمارية','Clarify'],['BOQ','5 بنود','جهد','Decompose']]
    for r,row in enumerate(vals):
        for c,v in enumerate(row): tbl.cell(r,c).text=v
    sl=prs.slides.add_slide(blank); cd=ChartData(); cd.categories=['Strategy','Team','Solution']; cd.add_series('Points',(20,20,30)); sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(1.2),Inches(1.4),Inches(10.5),Inches(4.8),cd)
    sl=prs.slides.add_slide(blank)
    for i,x in enumerate([.8,3.2,5.6,8,10.4]):
        sh=sl.shapes.add_shape(1,Inches(x),Inches(2.6),Inches(1.8),Inches(1)); sh.text=f'Step {i+1}'
    sl=prs.slides.add_slide(blank); a=sl.shapes.add_textbox(Inches(1),Inches(1.6),Inches(5),Inches(4)); a.text_frame.text='What attracts us\nStrategic fit\nLarge transformation'; b=sl.shapes.add_textbox(Inches(7),Inches(1.6),Inches(5),Inches(4)); b.text_frame.text='What blocks GO\nEvidence gaps\nCommercial unknowns'
    sl=prs.slides.add_slide(blank); tbl=sl.shapes.add_table(5,2,Inches(2),Inches(1.2),Inches(9),Inches(5)).table
    vals=[['Gate','Status'],['Experience','Unknown'],['Team','Unknown'],['Commercial','Open'],['Decision','HOLD']]
    for r,row in enumerate(vals):
        for c,v in enumerate(row): tbl.cell(r,c).text=v
    prs.save(path)

def bad_pptx(path):
    prs=Presentation(); blank=prs.slide_layouts[6]
    for i in range(10):
        sl=prs.slides.add_slide(blank)
        for j in range(10):
            sh=sl.shapes.add_shape(1,Inches(.7+(j%5)*2.4),Inches(1+(j//5)*2.3),Inches(2),Inches(1.4)); sh.text=f'Box {j}'
    prs.save(path)

def page(i,strategy,render_path):
    h=sha256_file(render_path)
    return {'page_id':f'P{i:02d}','selected_render_hash':h,'selected_strategy':strategy,'selected_candidate_id':f'H{(i%5)+1}','hypotheses':hyps(),'visual_concept_id':f'VC{i}','production_render_id':f'PR{i}','render_kind':'PRODUCTION_PAGE_RENDER','actual_pixel_review':review(h),'repair_required':False,'repair_history':[],'final_qa_round':1,'hero_metric_proven':strategy=='NUMBER_LED','structured_grid_rendered':False}
