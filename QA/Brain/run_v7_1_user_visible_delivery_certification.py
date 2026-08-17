#!/usr/bin/env python3
from __future__ import annotations
from pathlib import Path
import sys,json,tempfile,hashlib
ROOT=Path(__file__).resolve().parents[2]
BRAIN=ROOT/'Rashad/Brain/runtime'; sys.path.insert(0,str(BRAIN))
from brain.artifact_gate import guard_composer
from brain.actual_output_qa import evaluate_page_output,evaluate_deck_output
from brain.product_inspector import inspect_pptx,sha256_file
from brain.delivery_gate import validate_user_visible_delivery
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image,ImageDraw

OUT=ROOT/'QA/Certification'; OUT.mkdir(exist_ok=True)

def png(path,text):
    im=Image.new('RGB',(800,450),'white'); d=ImageDraw.Draw(im); d.text((40,40),text,fill='black'); d.rectangle((80,120,720,360),outline='black',width=3); im.save(path); return sha256_file(path)

def review(h,score=92,ind=True):
    dims=['message_clarity','five_second_comprehension','visual_form_fitness','simplicity','executive_hierarchy','evidence_legibility','artifact_usefulness','specificity_to_page','rtl_typography','brand_fidelity','production_quality']
    return {'status':'PASS','independent':ind,'review_id':'REV-'+h[:10],'actor_id':'QA-INDEPENDENT-PIXEL','actual_render_hash':h,'scores':{k:score for k in dims},'generic_layout_swap_test':'PASS','artifact_skeptic_test':'PASS','five_second_test':'PASS','hard_blockers':[]}

def hyps():
    ss=[('STATEMENT_LED','MINIMAL'),('NUMBER_LED','MINIMAL'),('EVIDENCE_LED','ANALYTICAL'),('COMPARISON_LED','ANALYTICAL'),('HYBRID_EXHIBIT','HYBRID')]
    return [{'id':f'H{i+1}','communication_strategy':s,'strategy_family':f,'structural_signature':f'{s}|SIG{i}'} for i,(s,f) in enumerate(ss)]

def good_pptx(path,imgpath):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    # 1 image-led
    sl=prs.slides.add_slide(blank); sl.shapes.add_picture(str(imgpath),Inches(.6),Inches(.7),width=Inches(6.1)); tb=sl.shapes.add_textbox(Inches(7.2),Inches(1.4),Inches(5.3),Inches(2)); tb.text_frame.text='مشروع الابتكار والثقافة الرقمية'
    # 2 statement
    sl=prs.slides.add_slide(blank); tb=sl.shapes.add_textbox(Inches(1),Inches(2.2),Inches(11),Inches(2)); p=tb.text_frame.paragraphs[0]; p.text='الفرصة جذابة، لكن قرار الدخول يحتاج إغلاق بوابات الجاهزية'; p.font.size=Pt(28)
    # 3 number
    sl=prs.slides.add_slide(blank); tb=sl.shapes.add_textbox(Inches(1),Inches(1.4),Inches(5),Inches(2.2)); p=tb.text_frame.paragraphs[0]; p.text='87.5%'; p.font.size=Pt(64); tb2=sl.shapes.add_textbox(Inches(6.3),Inches(2),Inches(5),Inches(2)); tb2.text_frame.text='المطلوب من بقية النقاط إذا فقدنا نقاط الخبرة'
    # 4 table
    sl=prs.slides.add_slide(blank); tbl=sl.shapes.add_table(4,4,Inches(.8),Inches(1.5),Inches(11.7),Inches(4.5)).table
    vals=[['Issue','Evidence','Impact','Action'],['مدة العقد','غير مكتملة','تسعير','Clarify'],['Open Source','تعارض','معمارية','Clarify'],['BOQ','5 بنود','جهد','Decompose']]
    for r,row in enumerate(vals):
        for c,v in enumerate(row): tbl.cell(r,c).text=v
    # 5 chart
    sl=prs.slides.add_slide(blank); cd=ChartData(); cd.categories=['Strategy','Team','Solution']; cd.add_series('Points',(20,20,30)); sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(1.2),Inches(1.4),Inches(10.5),Inches(4.8),cd)
    # 6 process
    sl=prs.slides.add_slide(blank)
    for i,x in enumerate([.8,3.2,5.6,8,10.4]):
        sh=sl.shapes.add_shape(1,Inches(x),Inches(2.6),Inches(1.8),Inches(1)); sh.text=f'Step {i+1}'
    # 7 comparison
    sl=prs.slides.add_slide(blank); a=sl.shapes.add_textbox(Inches(1),Inches(1.6),Inches(5),Inches(4)); a.text_frame.text='What attracts us\nStrategic fit\nLarge transformation'; b=sl.shapes.add_textbox(Inches(7),Inches(1.6),Inches(5),Inches(4)); b.text_frame.text='What blocks GO\nEvidence gaps\nCommercial unknowns'
    # 8 decision table
    sl=prs.slides.add_slide(blank); tbl=sl.shapes.add_table(5,2,Inches(2),Inches(1.2),Inches(9),Inches(5)).table
    vals=[['Gate','Status'],['Experience','Unknown'],['Team','Unknown'],['Commercial','Open'],['Decision','HOLD']]
    for r,row in enumerate(vals):
        for c,v in enumerate(row): tbl.cell(r,c).text=v
    prs.save(path)

def bad_pptx(path):
    prs=Presentation(); blank=prs.slide_layouts[6]
    for i in range(10):
        sl=prs.slides.add_slide(blank)
        for j in range(10):
            sh=sl.shapes.add_shape(1,Inches(.7+(j%5)*2.4),Inches(1+(j//5)*2.3),Inches(2),Inches(1.4)); sh.text=f'Box {j}'
    prs.save(path)

def page(i,strategy,render_path):
    h=sha256_file(render_path)
    return {'page_id':f'P{i:02d}','selected_render_hash':h,'selected_strategy':strategy,'selected_candidate_id':f'H{(i%5)+1}','hypotheses':hyps(),'visual_concept_id':f'VC{i}','production_render_id':f'PR{i}','render_kind':'PRODUCTION_PAGE_RENDER','actual_pixel_review':review(h),'repair_required':False,'repair_history':[],'final_qa_round':1,'hero_metric_proven':strategy=='NUMBER_LED','structured_grid_rendered':False}

def record(arr,name,ok,detail=''):
    arr.append({'name':name,'status':'PASS' if ok else 'FAIL','detail':detail})

def main():
    tests=[]
    with tempfile.TemporaryDirectory() as td:
        td=Path(td); hero=td/'hero.png'; png(hero,'HERO')
        ppt=td/'good.pptx'; good_pptx(ppt,hero)
        strategies=['IMAGE_LED','STATEMENT_LED','NUMBER_LED','TABLE_LED','CHART_LED','PROCESS_LED','COMPARISON_LED','DECISION_LED']
        pages=[]
        for i,s in enumerate(strategies,1):
            rp=td/f'p{i}.png'; png(rp,f'{i}-{s}'); pages.append(page(i,s,rp))
        montage=td/'montage.png'; png(montage,'MONTAGE')
        dh=sha256_file(ppt); mh=sha256_file(montage)
        deck_review={'status':'PASS','independent':True,'review_id':'DECK-1','actor_id':'QA-DECK-INDEPENDENT','deck_sha256':dh,'montage_sha256':mh,'scores':{k:92 for k in ['narrative_rhythm','visual_variety','executive_coherence','cross_page_specificity','density_rhythm','brand_consistency','rtl_consistency','overall_partner_grade']},'generic_deck_swap_test':'PASS','diagram_overuse_test':'PASS','hard_blockers':[]}
        product=inspect_pptx(ppt,pages)
        record(tests,'good_varied_pptx_product_inspection',product['status']=='PASS',str(product.get('blockers')))
        dq=evaluate_deck_output(pages,visibility='USER_VISIBLE_ARTIFACT_DRAFT',deck_review=deck_review,deck_sha256=dh,montage_sha256=mh,product_inspection=product)
        record(tests,'good_user_visible_actual_output_qa',dq['status']=='PASS',str(dq.get('blockers')))
        dossier={'schema':'RASHAD_USER_VISIBLE_ARTIFACT_DELIVERY_V1','classification':'USER_VISIBLE_ARTIFACT_DRAFT','output_file_sha256':dh,'montage_sha256':mh,'pages':pages,'artifact_brain_execution_status':'PASS','production_render_status':'PASS','actual_output_qa_closed_loop_status':'PASS','deck_pixel_review':deck_review,'framework_certification_substitute':False,'independent_judgment_status':'NOT_EXECUTED','parity_status':'NOT_EXECUTED','proof_index_status':'NOT_EXECUTED','release':{}}
        dg=validate_user_visible_delivery(dossier,ppt)
        record(tests,'good_exact_file_delivery_allowed',dg['status']=='DELIVERY_ALLOWED',str(dg.get('blockers')))

        # concept render leakage
        x=dict(pages[1]); x['render_kind']='COMMUNICATION_STRATEGY_CONCEPT_RENDER_V3'; x['actual_pixel_review']=review(x['selected_render_hash'])
        record(tests,'concept_render_blocked_user_visible',evaluate_page_output(x,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCKED')
        # missing pixel
        x=dict(pages[1]); x['actual_pixel_review']={'status':'NOT_EXECUTED'}
        record(tests,'missing_pixel_review_blocked',evaluate_page_output(x,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCKED')
        # low quality
        x=dict(pages[1]); x['actual_pixel_review']=review(x['selected_render_hash'],score=79)
        record(tests,'low_quality_pixel_blocked',evaluate_page_output(x,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCKED')
        # wrong hash
        x=dict(pages[1]); q=review(x['selected_render_hash']); q['actual_render_hash']='0'*64; x['actual_pixel_review']=q
        record(tests,'pixel_hash_mismatch_blocked',evaluate_page_output(x,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCKED')
        # producer review
        x=dict(pages[1]); x['actual_pixel_review']=review(x['selected_render_hash'],ind=False)
        record(tests,'non_independent_pixel_review_blocked',evaluate_page_output(x,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCKED')
        # generic swap
        x=dict(pages[1]); q=review(x['selected_render_hash']); q['generic_layout_swap_test']='FAIL'; x['actual_pixel_review']=q
        record(tests,'generic_layout_swap_failure_blocks',evaluate_page_output(x,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCKED')
        # repair missing
        x=dict(pages[1]); x['repair_required']=True; x['repair_history']=[]
        record(tests,'required_repair_history_blocked',evaluate_page_output(x,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCKED')
        # framework substitute
        bad=dict(dossier); bad['framework_certification_substitute']=True
        record(tests,'framework_certification_cannot_substitute',validate_user_visible_delivery(bad,ppt)['status']=='BLOCK_DELIVERY')
        # file hash mismatch
        bad=dict(dossier); bad['output_file_sha256']='f'*64
        record(tests,'exact_file_hash_mismatch_blocks',validate_user_visible_delivery(bad,ppt)['status']=='BLOCK_DELIVERY')
        # deck pixel mismatch
        bad=dict(dossier); bad['deck_pixel_review']=dict(deck_review); bad['deck_pixel_review']['deck_sha256']='e'*64
        record(tests,'deck_review_hash_mismatch_blocks',validate_user_visible_delivery(bad,ppt)['status']=='BLOCK_DELIVERY')
        # bad repetitive pptx
        bp=td/'bad.pptx'; bad_pptx(bp); insp=inspect_pptx(bp)
        record(tests,'shape_only_repetitive_deck_blocked',insp['status']=='BLOCKED' and ('PPTX_STRUCTURAL_MONOTONY' in insp['blockers'] or 'PPTX_SHAPE_ONLY_ANALYTICAL_DECK_OVERUSE' in insp['blockers']),str(insp['blockers']))
        # semantic mismatch claims
        mismatch=[dict(p) for p in pages]; mismatch[0]=dict(mismatch[0]); mismatch[0]['selected_strategy']='IMAGE_LED'
        # inspect a pptx without image on slide 1
        plain=td/'plain.pptx'; bad_pptx(plain); ip=inspect_pptx(plain,mismatch[:8])
        record(tests,'image_led_without_image_blocked',any('IMAGE_LED_WITHOUT_IMAGE' in b for b in ip['blockers']))
        mismatch2=[dict(p) for p in pages]; mismatch2[4]=dict(mismatch2[4]); mismatch2[4]['selected_strategy']='CHART_LED'
        record(tests,'chart_led_contract_present',product['slides'][4]['charts']>=1)
        # concept stage composer guard
        state={'content_status':'PASS','evidence_status':'PASS','page_contract':{'x':1},'cognitive_packet':{'x':1},'artifact_intent':{'x':1},'semantic_graph':{'x':1},'hypotheses':hyps(),'render_evidence':[{'actual_render_hash':str(i)*64} for i in range(1,6)]}
        record(tests,'internal_concept_draft_allowed',guard_composer(state,'INTERNAL_CONCEPT_DRAFT')['status']=='PASS')
        record(tests,'same_state_user_visible_blocked',guard_composer(state,'USER_VISIBLE_ARTIFACT_DRAFT')['status']=='BLOCK_RENDER')
        record(tests,'legacy_artifact_draft_alias_is_user_visible_and_blocked',guard_composer(state,'ARTIFACT_DRAFT')['status']=='BLOCK_RENDER')

    status='PASS' if all(x['status']=='PASS' for x in tests) else 'FAIL'
    out={'suite':'V7.1 User Visible Artifact Delivery Certification','status':status,'passed':sum(x['status']=='PASS' for x in tests),'total':len(tests),'tests':tests}
    (OUT/'V7_1_USER_VISIBLE_DELIVERY_CERTIFICATION.json').write_text(json.dumps(out,ensure_ascii=False,indent=2),encoding='utf8')
    print(json.dumps({'suite':out['suite'],'status':status,'passed':out['passed'],'total':out['total'],'failed':[x for x in tests if x['status']!='PASS']},ensure_ascii=False,indent=2))
    return 0 if status=='PASS' else 2
if __name__=='__main__': raise SystemExit(main())
